"""
Parte 1: Import, helpers, slides 1-20
Genera PowerPoint DETALLADO con TODO el progreso de GeoVision-CLIP Cali
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE = RGBColor(0x00, 0x9D, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x3D, 0x5A)
ACCENT_PURPLE = RGBColor(0xA8, 0x5D, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xE0, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x88)
TABLE_HDR = RGBColor(0x00, 0x6D, 0xC6)
TABLE_ALT = RGBColor(0x18, 0x28, 0x48)
TABLE_BASE = RGBColor(0x12, 0x1E, 0x3A)
SUBTLE = RGBColor(0x22, 0x33, 0x55)
CARD_BG = RGBColor(0x12, 0x1E, 0x3A)

BASE_DIR = r"C:\Users\nicor\universidad\analitica\proyectos\proyecto 3.1"
OUTPUT_FILE = os.path.join(BASE_DIR, "docs", "informe", "GeoVision_CLIP_Detallado.pptx")
ARQ_IMG = os.path.join(BASE_DIR, "notebooks", "Situacion_1", "Arquitectura", "arquitecturaSIT1.png")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

def bg(slide, c=DARK_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = c

def rect(slide, l, t, w, h, c):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = c
    s.line.fill.background()
    return s

def tb(slide, l, t, w, h, txt, sz=18, c=WHITE, b=False, al=PP_ALIGN.LEFT, fn="Calibri"):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = fn
    p.alignment = al
    return bx

def bullets(slide, l, t, w, h, items, sz=14, c=LIGHT_GRAY, bc=None):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = (bc + "  ") if bc else "  "
        p.text = prefix + item
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return bx

def header(slide, num, title, sub=None):
    rect(slide, 0, 0, W, Pt(5), ACCENT_BLUE)
    tb(slide, Inches(0.6), Inches(0.25), Inches(0.6), Inches(0.5),
       f"{num:02d}", sz=26, c=ACCENT_BLUE, b=True)
    tb(slide, Inches(1.3), Inches(0.25), Inches(10), Inches(0.5),
       title, sz=26, c=WHITE, b=True)
    if sub:
        tb(slide, Inches(1.3), Inches(0.65), Inches(10), Inches(0.35),
           sub, sz=13, c=MEDIUM_GRAY)
    rect(slide, Inches(0.6), Inches(1.05), Inches(12.1), Pt(2), SUBTLE)

def card(slide, l, t, w, h, c=CARD_BG):
    return rect(slide, l, t, w, h, c)

def kpi(slide, l, t, w, h, label, val, color=ACCENT_GREEN):
    card(slide, l, t, w, h)
    rect(slide, l, t, w, Pt(4), color)
    tb(slide, l+Inches(0.1), t+Inches(0.12), w-Inches(0.2), Inches(0.45),
       val, sz=22, c=color, b=True, al=PP_ALIGN.CENTER)
    tb(slide, l+Inches(0.1), t+Inches(0.55), w-Inches(0.2), Inches(0.45),
       label, sz=10, c=LIGHT_GRAY, al=PP_ALIGN.CENTER)

def tbl(slide, l, t, col_ws, headers, rows):
    for j, (hdr, cw) in enumerate(zip(headers, col_ws)):
        x = l + sum(col_ws[:j])
        rect(slide, x, t, cw, Inches(0.32), TABLE_HDR)
        tb(slide, x, t, cw, Inches(0.32), hdr, sz=10, c=WHITE, b=True, al=PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        y = t + Inches(0.32) + i * Inches(0.28)
        bg_c = TABLE_ALT if i % 2 == 0 else TABLE_BASE
        for j, (cell, cw) in enumerate(zip(row, col_ws)):
            x = l + sum(col_ws[:j])
            rect(slide, x, y, cw, Inches(0.28), bg_c)
            is_b = (i == len(rows)-1)
            tb(slide, x, y, cw, Inches(0.28), str(cell), sz=9,
               c=WHITE if is_b else LIGHT_GRAY, b=is_b, al=PP_ALIGN.CENTER)

def section(slide, l, t, w, txt, color=ACCENT_BLUE):
    tb(slide, l, t, w, Inches(0.35), txt, sz=16, c=color, b=True)

def multi(slide, l, t, w, h, lines, sz=12, c=LIGHT_GRAY):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.name = "Calibri"
        p.space_after = Pt(2)

print("Creando presentacion detallada...")

# === SLIDE 1: PORTADA ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, RGBColor(0x08,0x0E,0x1A))
rect(s, 0, 0, Inches(0.15), H, ACCENT_BLUE)
rect(s, Inches(1.5), Inches(1.8), Inches(10.3), Pt(3), ACCENT_BLUE)
rect(s, Inches(1.5), Inches(5.0), Inches(10.3), Pt(3), ACCENT_GREEN)
tb(s, Inches(1.5), Inches(2.0), Inches(10.5), Inches(1.0), "GeoVision-CLIP Cali", sz=46, c=WHITE, b=True)
tb(s, Inches(1.5), Inches(2.8), Inches(10.5), Inches(0.6),
   "Estimacion de Contaminacion Atmosferica en Puntos No Muestreados", sz=20, c=LIGHT_GRAY)
tb(s, Inches(1.5), Inches(3.4), Inches(10.5), Inches(0.5),
   "Deep Learning Multimodal + Estadistica Geoespacial Avanzada", sz=16, c=ACCENT_BLUE)
tb(s, Inches(1.5), Inches(5.3), Inches(10), Inches(1.2),
   "Analitica de Datos I  Tercer Corte  Unidad III\nUniversidad Autonoma de Occidente  Facultad de Ingenierias\nNicolas Pena Irurita  Samuel Patino Lucumi\nDocentes: Prof. Carlos Ferro  Prof. Cristian E. Garcia",
   sz=12, c=MEDIUM_GRAY)
rect(s, Inches(10.5), Inches(6.2), Inches(2.3), Inches(0.35), ACCENT_ORANGE)
tb(s, Inches(10.5), Inches(6.2), Inches(2.3), Inches(0.35),
   "INFORME DE PROGRESO", sz=9, c=WHITE, b=True, al=PP_ALIGN.CENTER)

# === SLIDE 2: AGENDA ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 0, "Agenda Detallada", "45+ diapositivas  Pipeline  EDAs  Dataset  Modelo  SAE  Frontend")
items = [
    ("01-03", "Introduccion", "Problema, arquitectura, stack tecnologico"),
    ("04-12", "Pipeline ETL", "config.py, exportar.py, convertir_zarr, validar, manifest, trazabilidad"),
    ("13-24", "EDA Exhaustivo", "S5P, S2, MODIS, ERA5, DAGMA, barrios, hallazgos, correlaciones"),
    ("25-31", "Dataset Situacion 2", "v1 vs v2, percentiles, clases, split, diagnostico, mejoras"),
    ("32-38", "Modelo GeoVision-CLIP + SAE", "Arquitectura, RemoteCLIP, SAE post-train, AFE/AFC, KPIs"),
    ("39-41", "Situacion 3 (Avance)", "Tensor ConvLSTM, covariables, embeddings"),
    ("42-44", "Frontend + Publicacion", "Next.js, HF Hub, entregables"),
    ("45", "Roadmap y Cierre", "Proximos pasos, logros"),
]
for i, (num, title, desc) in enumerate(items):
    y = Inches(1.4) + i * Inches(0.65)
    tb(s, Inches(0.8), y, Inches(1.0), Inches(0.3), num, sz=12, c=ACCENT_BLUE, b=True)
    tb(s, Inches(1.8), y, Inches(3.0), Inches(0.25), title, sz=14, c=WHITE, b=True)
    tb(s, Inches(1.8), y+Inches(0.24), Inches(10), Inches(0.2), desc, sz=10, c=MEDIUM_GRAY)
    if i < 7:
        rect(s, Inches(1.8), y+Inches(0.5), Inches(10.3), Pt(1), SUBTLE)

# === SLIDE 3: PROBLEMA ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 1, "El Problema", "Contaminacion atmosferica en Santiago de Cali")
section(s, Inches(0.8), Inches(1.4), Inches(5.5), "Contexto urbano y ambiental")
bullets(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.0), [
    "Cali: 3a ciudad de Colombia, ~2.8 M hab., 564 km2",
    "Fuentes: parque automotor, industria Yumbo-Acopi, quemas de cana",
    "Resolucion 2254/2017: niveles maximos permisibles NO2, SO2, O3",
    "Red DAGMA: solo 9 estaciones  cobertura espacial insuficiente",
], sz=13)
section(s, Inches(7), Inches(1.4), Inches(5.5), "La brecha satelital", ACCENT_ORANGE)
bullets(s, Inches(7), Inches(1.8), Inches(5.5), Inches(2.0), [
    "Sentinel-5P TROPOMI: gases traza, ~3.5x5.5 km  no ve barrios",
    "Sentinel-2 MSI: 10 m resolucion, pero NO mide gases directamente",
    "Hipotesis: fusionar S2 (contexto) + S5P (etiqueta) via CLIP",
    "Downscaling estadistico: de 5.5 km a 10 m de resolucion efectiva",
], sz=13, c=LIGHT_GRAY)
card(s, Inches(1), Inches(4.5), Inches(11.3), Inches(1.0))
rect(s, Inches(1), Inches(4.5), Pt(6), Inches(1.0), ACCENT_ORANGE)
tb(s, Inches(1.3), Inches(4.55), Inches(10.7), Inches(0.45),
   "Podemos estimar NO2, SO2 y O3 en cualquier (lat, lon) de Cali...", sz=17, c=WHITE, b=True)
tb(s, Inches(1.3), Inches(4.95), Inches(10.7), Inches(0.4),
   "...con incertidumbre cuantificada, usando solo satelites gratuitos y 9 estaciones de validacion?", sz=14, c=ACCENT_BLUE)

# === SLIDE 4: ARQUITECTURA ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 2, "Arquitectura del Proyecto", "3 situaciones secuenciales  Pesos: 20% + 20% + 30% + 10% + 10% + 10%")
sits = [
    ("SIT. 1  PANEL", "Panel espacio-temporal", "20%", ACCENT_BLUE,
     ["Panel >50 GB multi-fuente (7 fuentes)",
      "Pipeline ETL distribuido con Dask",
      "Zarr chunked + Parquet en GCS",
      "Manifest JSON con hashes MD5",
      "EDA exhaustivo (5 notebooks)",
      "COMPLETADO  66.78 GB"]),
    ("SIT. 2  CLIP+SAE", "Aprendizaje multimodal", "20%", ACCENT_GREEN,
     ["Dataset 2,276 pares imagen-texto",
      "RemoteCLIP ViT-B/32 + MiniLM",
      "Sparse Autoencoders (SAE)",
      "InfoNCE + regularizacion L1",
      "Validacion AFE + AFC psicometrica",
      "EN ENTRENAMIENTO"]),
    ("SIT. 3  KRIGING", "Prediccion geoestadistica", "30%", ACCENT_ORANGE,
     ["ConvLSTM bidireccional (8 fechas)",
      "Variograma espacio-temporal",
      "ST-Kriging  superficie continua + s2",
      "LOO-CV vs 9 estaciones DAGMA",
      "Moran I + LISA + mapas de clusters",
      "PENDIENTE"]),
]
for i, (title, subtitle, weight, color, items) in enumerate(sits):
    x = Inches(0.5) + i * Inches(4.2)
    card(s, x, Inches(1.4), Inches(3.95), Inches(5.2))
    rect(s, x, Inches(1.4), Inches(3.95), Pt(6), color)
    tb(s, x+Inches(0.15), Inches(1.55), Inches(3.6), Inches(0.35), title, sz=13, c=color, b=True)
    tb(s, x+Inches(0.15), Inches(1.85), Inches(3.6), Inches(0.5), subtitle, sz=10, c=LIGHT_GRAY)
    rect(s, x+Inches(2.6), Inches(1.55), Inches(1.1), Inches(0.28), color)
    tb(s, x+Inches(2.6), Inches(1.55), Inches(1.1), Inches(0.28), weight, sz=9, c=WHITE, b=True, al=PP_ALIGN.CENTER)
    bullets(s, x+Inches(0.15), Inches(2.5), Inches(3.6), Inches(3.5), items, sz=11, c=LIGHT_GRAY)

# === SLIDE 5: STACK ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 3, "Stack Tecnologico", "Herramientas por capa del proyecto")
stacks = [
    ("Almacenamiento", ["Google Cloud Storage (GCS)", "Zarr (cubos N-D chunked)", "Parquet particionado", "GeoTIFF (raw)"], ACCENT_BLUE),
    ("ETL Distribuido", ["Dask Distributed (4wx2t, 8GB)", "Xarray + xee (GEE bridge)", "Google Earth Engine API", "ThreadPoolExecutor (Zarr)"], ACCENT_GREEN),
    ("Deep Learning", ["PyTorch 2.1+", "PyTorch Lightning", "open-clip (RemoteCLIP)", "Transformers (MiniLM)", "LoRA (PEFT adapters)"], ACCENT_ORANGE),
    ("Geoestadistica", ["PyKrige (OK3D)", "PySAL (esda.Moran)", "GeoPandas", "scikit-learn", "factor-analyzer + semopy"], ACCENT_RED),
    ("Backend / API", ["FastAPI + Uvicorn", "Python 3.10", "Docker multi-stage", "HuggingFace Hub"], ACCENT_PURPLE),
    ("Frontend", ["Next.js 16 + React 19", "Leaflet / React-Leaflet", "TypeScript + Tailwind", "Zustand + React Query", "Modo oscuro"], ACCENT_CYAN),
]
for i, (cat, techs, color) in enumerate(stacks):
    col, row = i % 3, i // 3
    x = Inches(0.5) + col * Inches(4.2)
    y = Inches(1.4) + row * Inches(2.8)
    card(s, x, y, Inches(3.95), Inches(2.5))
    rect(s, x, y, Pt(5), Inches(2.5), color)
    tb(s, x+Inches(0.2), y+Inches(0.1), Inches(3.5), Inches(0.3), cat, sz=14, c=color, b=True)
    for j, tech in enumerate(techs):
        tb(s, x+Inches(0.2), y+Inches(0.45)+j*Inches(0.32), Inches(3.5), Inches(0.3),
           tech, sz=10, c=LIGHT_GRAY)

# === SLIDE 6: PIPELINE VISION GENERAL ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 4, "Pipeline ETL: Vision General", "11 modulos  6 fuentes satelitales  28 ejecuciones trazadas")
stages = [
    ("config.py", "Config central\nBBOX, fuentes, Dask", ACCENT_BLUE),
    ("exportar.py", "Dask: GEE  GCS\nMD5 inline", ACCENT_GREEN),
    ("convertir_zarr.py", "GeoTIFF  Zarr 4D\nThreadPoolExecutor", ACCENT_ORANGE),
    ("validar_zarr.py", "Integridad:\nshape, %NaN", ACCENT_RED),
    ("manifest.py", "Manifest global\nMD5 consolidado", ACCENT_PURPLE),
    ("generar_dataset\n_sit2_v2.py", "Pares img-txt\nSplit 70/15/15", ACCENT_CYAN),
]
for i, (name, desc, color) in enumerate(stages):
    x = Inches(0.3) + i * Inches(2.15)
    card(s, x, Inches(1.5), Inches(1.95), Inches(1.2))
    rect(s, x, Inches(1.5), Inches(1.95), Pt(4), color)
    tb(s, x+Inches(0.08), Inches(1.58), Inches(1.8), Inches(0.35), name, sz=10, c=color, b=True, al=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.08), Inches(1.9), Inches(1.8), Inches(0.65), desc, sz=8, c=LIGHT_GRAY, al=PP_ALIGN.CENTER)
    if i < 5:
        tb(s, x+Inches(1.95), Inches(1.8), Inches(0.25), Inches(0.4), ">", sz=18, c=MEDIUM_GRAY, b=True)
section(s, Inches(0.8), Inches(3.0), Inches(5.5), "Modulos de soporte")
bullets(s, Inches(0.8), Inches(3.4), Inches(5.5), Inches(2.5), [
    "trazabilidad/sistema.py  Run context manager, eventos JSONL",
    "trazabilidad/dask_plugin.py  forwarding de logs workers  cliente",
    "silenciar_warnings.py  suprime warnings de GDAL, GCloud, etc.",
    "publicar_hf.py  subir/descargar desde HuggingFace Hub",
    "preparar_covariables_convlstm.py  features para Sit. 3",
], sz=12)
section(s, Inches(7), Inches(3.0), Inches(5.5), "Ejecuciones registradas (runs/)", ACCENT_ORANGE)
bullets(s, Inches(7), Inches(3.4), Inches(5.5), Inches(3.0), [
    "28 directorios en runs/ con timestamp UTC",
    "Cada run: log.txt + eventos.jsonl + system.json",
    "Exportar: 14 ejecuciones (S2, S5P, ERA5, MODIS)",
    "Dataset Sit2 v1/v2: 10 ejecuciones",
    "Validar Zarr: 2 ejecuciones",
    "Publicar HF: 2 ejecuciones",
], sz=12, c=LIGHT_GRAY)

# === SLIDE 7: config.py ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 5, "config.py", "Configuracion central del pipeline (112 lineas)")
section(s, Inches(0.8), Inches(1.3), Inches(5.5), "Constantes globales")
multi(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.0), [
    "PROJECT_GEE = 'proyecto3ia-494900'",
    "PROJECT_GCP = 'proyecto3ia-494900'",
    "BUCKET = 'geovision-cali'",
    "BBOX = [-76.65, 3.30, -76.30, 3.65]",
    "  (Cali + Yumbo + Acopi, mas amplio que consigna)",
    "FECHA_INICIO = '2019-01-01'",
    "FECHA_FIN = '2023-12-31'",
    "DASK_N_WORKERS = 4, DASK_THREADS_PER_WORKER = 2",
    "DASK_MEMORY_LIMIT = '4GB'",
], sz=10, c=ACCENT_CYAN)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Diccionario FUENTES (6 fuentes)")
multi(s, Inches(7), Inches(1.7), Inches(5.8), Inches(4.5), [
    "1. S5P/L3_NO2: 3 bandas, escala 1113m",
    "2. S5P/L3_SO2: 2 bandas, escala 1113m",
    "3. S5P/L3_O3: 2 bandas, escala 1113m",
    "4. S2_SR_HARMONIZED: 13 bandas, 10m, por_banda",
    "   Dask override: 2 workers x 1 thread (S2 = 99%)",
    "5. ERA5/HOURLY: 7 bandas, escala 27830m",
    "6. MODIS/MCD19A2: 4 bandas, escala 927m",
    "",
    "Cada fuente: bandas, escala, prefijo GCS, modo",
    "  modo multibanda: 1 GeoTIFF por imagen",
    "  modo por_banda: 1 GeoTIFF por (imagen, banda)",
    "    S2  genera 13 archivos por escena!",
    "",
    "get_fuente(): normaliza alias (no2, s2, era5...)",
], sz=10, c=LIGHT_GRAY)

# === SLIDE 8: exportar.py ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 6, "exportar.py", "ETL distribuido: GEE  GCS (794 lineas)")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Arquitectura de descarga")
bullets(s, Inches(0.8), Inches(1.7), Inches(5.8), Inches(3.5), [
    "Clase RegistroArchivo: fuente, img_id, banda, path, md5, size, bbox",
    "Lista imagenes GEE  construye tareas Dask",
    "Cache de 2 niveles:",
    "  Nivel 1: indice prelistado de GCS",
    "  Nivel 2: blob.exists() inline en worker",
    "Broadcast del cache_index a workers (scatter)",
    "Backoff exponencial con jitter para HTTP (429/5xx)",
    "Sidecar MD5: hash calculado durante transmision",
    "Manifest parcial por fuente (manifest_partial.json)",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Funciones clave", ACCENT_ORANGE)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(4.5), [
    "extraer_fecha(img_id, fuente_id):",
    "  Parseo por fuente:",
    "  S5P: YYYYMMDDTHHMMSS",
    "  S2: YYYYMMDD",
    "  MODIS: YYYYDDD (dia juliano)",
    "  ERA5: YYYYMMDDTHH",
    "",
    "_descargar_con_reintento(url, timeout=300s):",
    "  Backoff: 2s, 4s, 8s... hasta 60s",
    "  Jitter +/-25%",
    "  Max 5 reintentos en 429/5xx",
    "",
    "construir_tareas():",
    "  multibanda: 1 tarea por imagen",
    "  por_banda: 1 tarea por (imagen, banda)",
    "  S2: 13 tareas por escena!",
], sz=10, c=LIGHT_GRAY)

# === SLIDE 9: CACHING ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 7, "exportar.py: Estrategia de Caching", "Evita redes cargar 283,437 archivos")
steps = [
    ("1. Prelistar GCS", "Lee prefijo/raw/*.tif\nConstruye indice:\n{nombre: (size, md5)}", ACCENT_BLUE),
    ("2. Broadcast", "Scatter del indice\na todos los workers\nDask", ACCENT_GREEN),
    ("3. Cache hit?", "Worker consulta\nindice local\nExiste? MD5 ok?", ACCENT_ORANGE),
    ("4. Si cache hit", "Registro cache\nSin descarga\nSin costo GEE", ACCENT_GREEN),
    ("5. Si cache miss", "Descarga URL GEE\nMD5 inline\nSubida a GCS", ACCENT_RED),
    ("6. Manifest", "Registro descargado\nMD5 verificado\nflush a JSON", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.2) + i * Inches(2.2)
    card(s, x, Inches(1.7), Inches(2.0), Inches(1.6))
    rect(s, x, Inches(1.7), Inches(2.0), Pt(4), color)
    tb(s, x+Inches(0.08), Inches(1.78), Inches(1.85), Inches(0.3), title, sz=10, c=color, b=True, al=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.08), Inches(2.1), Inches(1.85), Inches(1.0), desc, sz=8, c=LIGHT_GRAY, al=PP_ALIGN.CENTER)
    if i < 5:
        tb(s, x+Inches(2.0), Inches(2.2), Inches(0.25), Inches(0.4), ">", sz=16, c=MEDIUM_GRAY, b=True)
section(s, Inches(0.8), Inches(3.6), Inches(11), "Estadisticas de exportacion")
kpi(s, Inches(0.8), Inches(4.0), Inches(2.5), Inches(1.0), "Archivos exportados", "283,437", ACCENT_BLUE)
kpi(s, Inches(3.5), Inches(4.0), Inches(2.5), Inches(1.0), "Peso total", "66.78 GB", ACCENT_GREEN)
kpi(s, Inches(6.2), Inches(4.0), Inches(2.5), Inches(1.0), "Ejecuciones export", "14", ACCENT_ORANGE)
kpi(s, Inches(8.9), Inches(4.0), Inches(2.5), Inches(1.0), "S2 domina peso", "99.45%", ACCENT_RED)
bullets(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.5), [
    "S2 modo por_banda  13 tareas por escena  19,121 archivos generados",
    "MODIS 140,120 granulos: mas archivos pero mas pequenos",
    "Limite GEE 30 GB/dia  chunking mensual + distribucion temporal",
], sz=11, c=LIGHT_GRAY)

# === SLIDE 10: convertir_zarr.py ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 8, "convertir_zarr.py", "GeoTIFF  Cubos Zarr 4D (425 lineas)")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Estrategia: esqueleto + region writes")
bullets(s, Inches(0.8), Inches(1.7), Inches(5.8), Inches(3.0), [
    "Paso 1: Leer GeoTIFF muestra para shape espacial",
    "Paso 2: Crear Zarr vacio (solo metadatos)",
    "  xr.Dataset + darr.zeros + to_zarr(compute=False)",
    "Paso 3: Procesar batches con ThreadPoolExecutor",
    "  Cada hilo lee su item (GeoTIFF  ndarray float32)",
    "  Escribe region temporal con to_zarr(region=...)",
    "Paso 4: Consolidar metadatos (zarr.consolidate_metadata)",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Detalles tecnicos", ACCENT_ORANGE)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(4.0), [
    "Chunking adaptativo:",
    "  chunk_time = max(1, 268435456 // (b*H*W*4))",
    "  Cada chunk ~256 MB en RAM",
    "  Optimizado para lecturas parciales",
    "",
    "Thread-local storage:",
    "  _bucket_thread_local() reutiliza",
    "  storage.Client y Bucket por hilo",
    "",
    "Modos de agrupacion:",
    "  multibanda: 1 GeoTIFF = 1 item",
    "  por_banda: N GeoTIFFs (misma img_id)",
    "    se stackean en orden de bandas",
    "",
    "ThreadPoolExecutor (NO Dask):",
    "  Dask overhead alto para IO intensivo",
    "  8 workers default",
], sz=10, c=LIGHT_GRAY)

# === SLIDE 11: manifest + validar ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 9, "manifest.py + validar_zarr.py", "Gobernanza y aseguramiento de calidad")
section(s, Inches(0.8), Inches(1.3), Inches(5.8), "manifest.py (180 lineas)")
bullets(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.5), [
    "Consolida manifests parciales de todas las fuentes",
    "Descarga cada manifest_partial.json desde GCS",
    "Campos por archivo: ruta, MD5, size_bytes, fecha, bbox, escala",
    "Ordenamiento: (fuente, img_id, banda)",
    "Reporta fuentes faltantes en campo faltantes[]",
    "Estadisticas: n_archivos, size_bytes_total por fuente",
    "CLI: --subir-a-gcs, --salida-local",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "validar_zarr.py (180 lineas)", ACCENT_GREEN)
bullets(s, Inches(7), Inches(1.7), Inches(5.5), Inches(2.5), [
    "Verifica por fuente:",
    "  Zarr existe en GCS (consolidated=True)",
    "  Dimensiones: time, band, y, x",
    "  Bandas coinciden con configuracion",
    "  Rango temporal (min/max)",
    "  Chunks y dtype",
    "  %NaN muestreado (4 pasos equiespaciados)",
    "Exit code 0 si todo OK, 1 si corrupta",
    "Modo --json para integracion automatizada",
], sz=11, c=LIGHT_GRAY)
h = ["Fuente", "Shape", "Bandas", "Tiempos", "%NaN est.", "Estado"]
r = [
    ["S5P NO2", "(28,303, 3, 36, 36)", "3/3 ok", "2019-2023", "<1%", "OK"],
    ["S5P SO2", "(26,280, 2, 36, 36)", "2/2 ok", "2019-2023", "<1%", "OK"],
    ["S5P O3", "(25,798, 2, 36, 36)", "2/2 ok", "2019-2023", "<1%", "OK"],
    ["ERA5", "(43,815, 7, 5, 5)", "7/7 ok", "2019-2023", "<5%", "OK"],
    ["MODIS", "(140,120, 4, 5, 5)", "4/4 ok", "2019-2023", "<10%", "OK"],
]
tbl(s, Inches(0.5), Inches(4.6), [Inches(1.8), Inches(2.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.2)], h, r)

# === SLIDE 12: TRAZABILIDAD ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 10, "Sistema de Trazabilidad", "trazabilidad/sistema.py + dask_plugin.py")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Clase Run (209 lineas)")
bullets(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.5), [
    "Context manager que representa una corrida del pipeline",
    "Crea: runs/<YYYYMMDDTHHMMSS>_<nombre>/",
    "  log.txt  logging humano con timestamps",
    "  eventos.jsonl  JSON estructurado (1 por linea)",
    "",
    "Metodos:",
    "  Run(nombre, contexto)  init + ContextVar",
    "  .evento(nombre, **kw)  evento estructurado",
    "  .info(), .warning(), .debug(), .error()",
    "  .medir(nombre)  context manager con duracion",
    "",
    "Usa contextvars.ContextVar para hilos/async",
    "Sin dependencias externas (solo stdlib)",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Plugin Dask (60 lineas)", ACCENT_ORANGE)
bullets(s, Inches(7), Inches(1.7), Inches(5.5), Inches(2.0), [
    "forward_logs(client):",
    "  Registra callback en setup de workers",
    "  Cada worker: logging a stdout con [worker]",
    "  Silencia: distributed, urllib3, googleapiclient",
    "3 niveles de fallback defensivo",
    "  register_worker_callbacks",
    "  client.run()",
    "  warning sin romper pipeline",
], sz=11, c=LIGHT_GRAY)
section(s, Inches(7), Inches(3.5), Inches(5.5), "Eventos registrados", ACCENT_GREEN)
multi(s, Inches(7), Inches(3.9), Inches(5.5), Inches(2.5), [
    "Ejemplo eventos en eventos.jsonl:",
    '  {"evento": "run_inicio", "ts": "...", "nombre": "exportar"}',
    '  {"evento": "tarea_fin", "estado": "cache", "md5": "..."}',
    '  {"evento": "medir", "duracion_s": 123.4}',
    '  {"evento": "run_fin", "duracion_s": 3600, "exito": true}',
    "",
    "28 ejecuciones registradas  historial completo",
], sz=10, c=LIGHT_GRAY)

# === SLIDE 13: MODULOS AUXILIARES ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 11, "Modulos Auxiliares", "publicar_hf.py + silenciar_warnings.py + preparar_covariables_convlstm.py")
section(s, Inches(0.8), Inches(1.3), Inches(5.8), "publicar_hf.py (333 lineas)")
bullets(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.0), [
    "Sube/descarga dataset desde HuggingFace Hub",
    "Excluye siempre TIFFs crudos (**/raw/**)",
    "Modo ligero: excluye S2 (~66 GB)",
    "Modo ambiental: solo ERA5 + MODIS",
    "Usa upload_folder de huggingface_hub",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "silenciar_warnings.py (105 lineas)", ACCENT_GREEN)
bullets(s, Inches(7), Inches(1.7), Inches(5.5), Inches(2.0), [
    "Suprime warnings de:",
    "  google.auth._default (credenciales Cloud)",
    "  google.api_core (FutureWarning Python 3.10)",
    "  pkg_resources (DeprecationWarning)",
    "  GDAL/libtiff (CE_Warning SamplesPerPixel)",
    "Configura PYTHONWARNINGS para procesos hijos",
], sz=11, c=LIGHT_GRAY)
section(s, Inches(0.8), Inches(3.5), Inches(11), "preparar_covariables_convlstm.py (366 lineas)")
bullets(s, Inches(0.8), Inches(3.9), Inches(11.5), Inches(3.0), [
    "Carga metadatos.parquet + percentiles.json",
    "Features: doy_sin/cos, mes_sin/cos (circular encoding)",
    "Posicion z-score: lat_norm, lon_norm",
    "Contaminantes normalizados: no2_norm, so2_norm, o3_norm (p50/p90)",
    "Conversion mol/m2  ug/m3 (asumiendo 8 km troposfera Cali)",
    "Gap temporal desde tile anterior en misma celda 0.05 grados",
    "Opcional: ERA5 (T2m, BLH, viento, presion) y MODIS AOD",
    "Salida: covariables.parquet con ~20 columnas",
], sz=11, c=LIGHT_GRAY)

# === SLIDE 14: MANIFEST GLOBAL ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 12, "Manifest Global", "66.78 GB  283,437 archivos  MD5 verificables")
h = ["Fuente", "Archivos", "Tamano", "% Peso"]
r = [
    ["Sentinel-2 SR Harmonized", "19,121", "66.41 GB", "99.45%"],
    ["S5P NO2", "28,303", "0.08 GB", "0.11%"],
    ["S5P SO2", "26,280", "0.03 GB", "0.05%"],
    ["S5P O3", "25,798", "0.06 GB", "0.08%"],
    ["MODIS MCD19A2", "140,120", "0.11 GB", "0.16%"],
    ["ERA5", "43,815", "0.10 GB", "0.15%"],
    ["TOTAL", "283,437", "66.78 GB", "100%"],
]
tbl(s, Inches(0.8), Inches(1.5), [Inches(3.0), Inches(2.0), Inches(2.0), Inches(1.5)], h, r)
section(s, Inches(0.8), Inches(4.0), Inches(11), "Metadatos globales del manifest v1.0")
multi(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.5), [
    "Bucket: gs://geovision-cali (proyecto: proyecto3ia-494900)",
    "BBox: [-76.65, 3.30, -76.30, 3.65] (WGS84)",
    "Ventana temporal: 2019-01-01 a 2023-12-31",
    "Umbral 50 GB: CUMPLIDO (66.78 GB = +33.6% margen)",
    "Reproducibilidad: Docker (python:3.11-slim multi-stage)",
    "  requirements.txt congelado, seed fija 42",
    "",
    "Cada archivo en manifest:",
    "  fuente, img_id, banda, path GCS, MD5 hex, size_bytes,",
    "  fecha_adquisicion ISO, bbox, escala_m, estado",
], sz=12, c=LIGHT_GRAY)

# === SLIDE 15: EDA VISION GENERAL ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 13, "EDA: Vision General", "5 notebooks  30+ visualizaciones  7 fuentes de datos")
kpi(s, Inches(0.5), Inches(1.4), Inches(2.5), Inches(1.0), "Notebooks EDA", "5", ACCENT_BLUE)
kpi(s, Inches(3.2), Inches(1.4), Inches(2.5), Inches(1.0), "Visualizaciones", "30+", ACCENT_GREEN)
kpi(s, Inches(5.9), Inches(1.4), Inches(2.5), Inches(1.0), "Fuentes analizadas", "7", ACCENT_ORANGE)
kpi(s, Inches(8.6), Inches(1.4), Inches(2.5), Inches(1.0), "Hallazgos clave", "8+", ACCENT_PURPLE)
kpi(s, Inches(11.3), Inches(1.4), Inches(1.5), Inches(1.0), "CSV generados", "3", ACCENT_CYAN)
h = ["Notebook", "Fuente", "Variables", "Viz", "Hallazgos clave"]
r = [
    ["EDA_Sentinel5P.ipynb", "S5P NO2/SO2/O3", "3 gases, 28k escenas", "12", "Cobertura 7-53%, estacionalidad"],
    ["EDA_Sentinel2.ipynb", "S2 MSI L2A", "13 bandas, indices", "10", "Nubosidad 67%, separabilidad"],
    ["MODIS_MCD19A2_.ipynb", "MODIS AOD", "4 bandas AOD", "6", "Hotspot Yumbo, AOD ~0.28"],
    ["ERA5_.ipynb", "ERA5-Land", "T2m, viento, BLH, HR", "14", "BLH vs NO2 r=-0.52"],
    ["EDA_DAGMA.ipynb", "DAGMA SVCASC", "33,580 reg, 87 cols", "10", "NO2 solo Univalle"],
]
tbl(s, Inches(0.4), Inches(2.7), [Inches(2.2), Inches(1.8), Inches(2.0), Inches(0.6), Inches(4.8)], h, r)
section(s, Inches(0.8), Inches(5.0), Inches(11), "Objetivos del EDA")
bullets(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.5), [
    "Caracterizar cobertura temporal real (impacto nubes tropicales)",
    "Calibrar percentiles S5P sobre Cali para etiquetado semi-supervisado",
    "Detectar huecos y limitaciones del ground truth DAGMA",
    "Identificar correlaciones entre contaminantes y meteorologia",
    "Justificar decisiones de filtrado (SCL, umbrales de nubosidad)",
], sz=12, c=LIGHT_GRAY)

# === SLIDE 16: EDA S5P ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 14, "EDA: Sentinel-5P TROPOMI", "28,303 escenas NO2  26,280 SO2  25,798 O3 (2019-2023)")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Cobertura y disponibilidad")
multi(s, Inches(0.8), Inches(1.7), Inches(5.8), Inches(4.5), [
    "Colecciones: COPERNICUS/S5P/OFFL/L3_{NO2,SO2,O3}",
    "Zarr: gs://geovision-cali/Sentinel5P/{NO2,SO2,O3}/panel.zarr",
    "",
    "Cobertura efectiva (pixeles validos):",
    "  NO2: ~25%  impacto severo de nubes tropicales",
    "  SO2: ~53%  mejor cobertura entre los 3",
    "  O3:  ~7%  muy sensible a nubosidad en tropicos",
    "",
    "Patron orbital: ~25-45 escenas/mes (regular)",
    "Parseo system:index  YYYYMMDDTHHMMSS",
], sz=11, c=LIGHT_GRAY)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Series temporales", ACCENT_GREEN)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(4.5), [
    "NO2: picos en jul-ago (quemas de cana de azucar)",
    "  Caida 18% en 2020 (COVID-19 lockdown)",
    "  Recuperacion al 92% en 2023",
    "",
    "SO2: correlacion con NO2 (r=0.38)",
    "  Fuentes compartidas: transporte + industria",
    "",
    "O3: anticorrelacion con NO2 (r=-0.42)",
    "  Titulacion fotoquimica: NO + O3  NO2 + O2",
    "  Tipico de ambientes urbanos con NOx",
    "",
    "CV espacial: SO2 ~28.8% (alto), O3 ~0.35% (bajo)",
], sz=11, c=LIGHT_GRAY)

# === SLIDE 17: EDA S2 ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 15, "EDA: Sentinel-2 MSI", "13 bandas  10/20/60 m  736 escenas (2019-2023)")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Nubosidad y disponibilidad")
multi(s, Inches(0.8), Inches(1.7), Inches(5.8), Inches(2.5), [
    "Coleccion: COPERNICUS/S2_SR_HARMONIZED",
    "736 escenas (S2A+S2B) en 5 anos, 6 tiles",
    "Solo 242 (32.9%) con <60% de nubosidad",
    "Promedio: 4.2 escenas utiles/mes (max ~12)",
    "Minimos en abril-mayo y octubre-noviembre (lluvias)",
    "Mejores meses: junio-agosto (temporada seca)",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Indices espectrales", ACCENT_GREEN)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(2.5), [
    "NDVI = (B8-B4)/(B8+B4)  vegetacion",
    "  Yumbo Industrial: 0.1-0.2",
    "  Cali Urbano: 0.3-0.5",
    "  Pance Rural: 0.6-0.8",
    "BSI = ((B11+B4)-(B8+B2))/((B11+B4)+(B8+B2))",
    "  suelo desnudo, corredor Yumbo-Acopi",
    "NDBI = (B11-B8)/(B11+B8)  superficie construida",
], sz=11, c=LIGHT_GRAY)
section(s, Inches(0.8), Inches(4.0), Inches(11), "Perfiles espectrales y separabilidad de clases")
bullets(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.5), [
    "4 clases de cobertura: vegetacion densa, suelo desnudo, agua, zona construida",
    "Separacion 0.12-0.18 uds reflectancia en banda 8 (NIR, 842 nm)",
    "SWIR (B11-B12) discrimina suelo desnudo de zona construida",
    "3 bloques alta correlacion: Visible (B2-B4, r>0.98), RedEdge+NIR, SWIR",
    "Fundamenta S2 como fuente de covariables proxy para downscaling",
], sz=11, c=LIGHT_GRAY)

# === SLIDE 18: EDA MODIS ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 16, "EDA: MODIS MCD19A2 (AOD)", "140,120 granulos  Profundidad optica de aerosoles  2019-2023")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Parametros")
bullets(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.0), [
    "Coleccion: MODIS/061/MCD19A2_GRANULES (MAIAC)",
    "Bandas: Optical_Depth_047, _055, AOD_Uncertainty",
    "Escala nativa: 927 m (1 km), diaria",
    "Factor de escala: x0.001 para valores fisicos",
    "152,458 granulos totales sobre Cali",
    "Proxy indirecto de PM2.5 y PM10",
    "Cuota GEE tier no comercial alcanzada",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Hallazgos", ACCENT_ORANGE)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(3.0), [
    "AOD 550 nm promedio: ~0.28 sobre Cali",
    "Maximos estacionales: ~0.45 en jul-ago",
    "  Coincide con temporada de quemas de cana",
    "",
    "Hotspot: Corredor Yumbo-Acopi",
    "  AOD ~0.12 superior al promedio urbano",
    "  Consistente con fuentes industriales",
    "",
    "AOD_047 sistematicamente mayor que AOD_055",
    "Gradiente espacial: Oeste bajo (0.1-0.2)",
    "  Centro/Este alto (0.3-0.5)",
], sz=11, c=LIGHT_GRAY)

# === SLIDE 19: EDA ERA5 ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 17, "EDA: ERA5-Land (Meteorologia)", "43,815 archivos horarios  292 bandas  2019-2023")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Variables seleccionadas")
bullets(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.5), [
    "temperature_2m (T2m)  C (K-273.15)",
    "boundary_layer_height (BLH) en metros",
    "u_component_of_wind_10m, v_component_of_wind_10m",
    "total_precipitation  mm (x1000)",
    "relative_humidity % (formula Magnus)",
    "Escala: ~0.25 grados (~28 km)",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Correlaciones clave", ACCENT_GREEN)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(2.5), [
    "BLH vs NO2: r = -0.52  LA MAS PREDICTIVA",
    "  BLH baja (inversion termica) atrapa contaminantes",
    "",
    "T2m vs NO2: r = -0.38",
    "  Temperaturas altas aceleran reacciones",
    "",
    "HR vs O3: r = 0.41",
    "  Mayor humedad  mas ozono",
    "",
    "Viento vs SO2: r = -0.29",
    "  Vientos fuertes dispersan SO2 industrial",
], sz=11, c=LIGHT_GRAY)
section(s, Inches(0.8), Inches(4.5), Inches(11), "Climatologia de Cali (2019-2023)")
bullets(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(2.0), [
    "Temperatura media: 22-24 C, muy estable interanualmente",
    "BLH: 240-300 m, correlacion positiva fuerte con T2m (r~0.8)",
    "Humedad relativa: 80-85% promedio, alta todo el ano",
    "Viento predominante: SW (~285 grados), baja velocidad (~0.7 m/s)",
    "Precipitacion estacional: abril-mayo y octubre-noviembre",
], sz=11, c=LIGHT_GRAY)

# === SLIDE 20: EDA DAGMA ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, 18, "EDA: Ground Truth DAGMA/SVCASC", "33,580 registros  9 estaciones  2020-2023")
h = ["Estacion", "Lat", "Lon", "Contaminantes disponibles"]
r = [
    ["Base Aerea", "3.4646", "-76.5142", "O3, PM10, meteorologia"],
    ["Canaveralejo", "3.4189", "-76.5417", "SO2, O3, PM, meteo"],
    ["Compartir", "3.4306", "-76.4764", "O3, PM10, meteo"],
    ["La Ermita", "3.4515", "-76.5322", "SO2, O3, PM, meteo"],
    ["La Flora", "3.4918", "-76.5142", "SO2, O3, PM, meteo"],
    ["ERA Obrero", "3.4537", "-76.5208", "O3, PM, meteo"],
    ["Pance", "3.3278", "-76.5486", "O3, PM10, meteo"],
    ["Transitoria Navarro", "3.4789", "-76.4892", "O3, PM10, meteo"],
    ["Univalle", "3.3779", "-76.5338", "NO2, SO2, O3, PM, meteo"],
]
tbl(s, Inches(0.4), Inches(1.5), [Inches(2.2), Inches(0.8), Inches(0.8), Inches(4.0)], h, r)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Cobertura por variable", ACCENT_ORANGE)
multi(s, Inches(7.2), Inches(1.5), Inches(5.3), Inches(3.0), [
    "NO2: SOLO en Univalle (69% cobertura)  CRITICO",
    "O3: 6 estaciones (48-83%)",
    "SO2: 5 estaciones (9-88%)",
    "PM2.5: 3 estaciones (62-83%)",
    "PM10: 7 estaciones (80-89%)  mejor cubierto",
    "Black Carbon: 3% (<5%  excluido)",
    "UV-PM: 3% (<5%  excluido)",
    "Meteo (temp, HR, viento): 6-9 estaciones",
], sz=10, c=LIGHT_GRAY)
section(s, Inches(0.8), Inches(5.0), Inches(11), "Hallazgos principales")
bullets(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.5), [
    "Correlacion NO2-O3: r = -0.42 (titulacion fotoquimica, consistente con S5P)",
    "O3-temperatura: r ~ 0.6 (formacion fotoquimica de ozono)",
    "Efecto COVID-19: reduccion NO2 y PM en 2020, recuperacion 2021-2023",
    "Episodios extremos: SO2 hasta 479 ug/m3 en Ermita, PM10 hasta 225.6 ug/m3",
    "Ciclo diurno y estacional: O3 pico mitad de ano, PM en meses secos (Dic-Mar)",
], sz=11, c=LIGHT_GRAY)


print(f"Total slides p1: {len(prs.slides)}")
prs.save(OUTPUT_FILE)
print(f"Guardado: {OUTPUT_FILE}")

