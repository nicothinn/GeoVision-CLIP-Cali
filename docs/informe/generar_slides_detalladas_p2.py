"""
Parte 2: Slides 21-45 (EDA Barrios, Dataset, Modelo, SAE, AFE/AFC, Frontend, Cierre)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import sys

# Cargar presentacion existente
BASE_DIR = r"C:\Users\nicor\universidad\analitica\proyectos\proyecto 3.1"
OUTPUT_FILE = os.path.join(BASE_DIR, "docs", "informe", "GeoVision_CLIP_Detallado.pptx")

from pptx import Presentation as P
prs = P(OUTPUT_FILE)
W, H = prs.slide_width, prs.slide_height

# Re-definir helpers (no podemos importar de p1)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE = RGBColor(0x00, 0x9D, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x3D, 0x5A)
ACCENT_PURPLE = RGBColor(0xA8, 0x5D, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xE0, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x88)
TABLE_HDR = RGBColor(0x00, 0x6D, 0xC6)
TABLE_ALT = RGBColor(0x18, 0x28, 0x48)
TABLE_BASE = RGBColor(0x12, 0x1E, 0x3A)
SUBTLE = RGBColor(0x22, 0x33, 0x55)
CARD_BG = RGBColor(0x12, 0x1E, 0x3A)

def bg(slide, c=DARK_BG):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = c

def rect(slide, l, t, w, h, cc):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = cc; s.line.fill.background()
    return s

def tb(slide, l, t, w, h, txt, sz=18, cc=WHITE, b=False, al=PP_ALIGN.LEFT, fn="Calibri"):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz)
    p.font.color.rgb = cc; p.font.bold = b; p.font.name = fn; p.alignment = al
    return bx

def bullets(slide, l, t, w, h, items, sz=14, cc=LIGHT_GRAY, bc=None):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = (bc + "  ") if bc else "  "
        p.text = prefix + item; p.font.size = Pt(sz); p.font.color.rgb = cc
        p.font.name = "Calibri"; p.space_after = Pt(4)
    return bx

def header(slide, num, title, sub=None):
    rect(slide, 0, 0, W, Pt(5), ACCENT_BLUE)
    tb(slide, Inches(0.6), Inches(0.25), Inches(0.6), Inches(0.5),
       f"{num:02d}", sz=26, cc=ACCENT_BLUE, b=True)
    tb(slide, Inches(1.3), Inches(0.25), Inches(10), Inches(0.5), title, sz=26, cc=WHITE, b=True)
    if sub:
        tb(slide, Inches(1.3), Inches(0.65), Inches(10), Inches(0.35), sub, sz=13, cc=MEDIUM_GRAY)
    rect(slide, Inches(0.6), Inches(1.05), Inches(12.1), Pt(2), SUBTLE)

def card(slide, l, t, w, h, cc=CARD_BG):
    return rect(slide, l, t, w, h, cc)

def kpi(slide, l, t, w, h, label, val, color=ACCENT_GREEN):
    card(slide, l, t, w, h)
    rect(slide, l, t, w, Pt(4), color)
    tb(slide, l+Inches(0.1), t+Inches(0.12), w-Inches(0.2), Inches(0.45),
       val, sz=22, cc=color, b=True, al=PP_ALIGN.CENTER)
    tb(slide, l+Inches(0.1), t+Inches(0.55), w-Inches(0.2), Inches(0.45),
       label, sz=10, cc=LIGHT_GRAY, al=PP_ALIGN.CENTER)

def tbl(slide, l, t, col_ws, headers, rows):
    for j, (hdr, cw) in enumerate(zip(headers, col_ws)):
        x = l + sum(col_ws[:j])
        rect(slide, x, t, cw, Inches(0.32), TABLE_HDR)
        tb(slide, x, t, cw, Inches(0.32), hdr, sz=10, cc=WHITE, b=True, al=PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        y = t + Inches(0.32) + i * Inches(0.28)
        bg_c = TABLE_ALT if i % 2 == 0 else TABLE_BASE
        for j, (cell, cw) in enumerate(zip(row, col_ws)):
            x = l + sum(col_ws[:j])
            rect(slide, x, y, cw, Inches(0.28), bg_c)
            is_b = (i == len(rows)-1)
            tb(slide, x, y, cw, Inches(0.28), str(cell), sz=9,
               cc=WHITE if is_b else LIGHT_GRAY, b=is_b, al=PP_ALIGN.CENTER)

def section(slide, l, t, w, txt, color=ACCENT_BLUE):
    tb(slide, l, t, w, Inches(0.35), txt, sz=16, cc=color, b=True)

def multi(slide, l, t, w, h, lines, sz=12, cc=LIGHT_GRAY):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = cc
        p.font.name = "Calibri"; p.space_after = Pt(2)

print("Agregando slides 21-45...")

# === SLIDE 21: EDA BARRIOS ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 19, "EDA: Barrios de Cali", "339 barrios  22 comunas  Shapefile oficial")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Datos geograficos")
bullets(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.0), [
    "Fuente: Shapefile oficial mc_barrios.shp",
    "CRS original: EPSG:6249 (MAGNA-SIRGAS / Cali-Valle 2009)",
    "Reproyectado a EPSG:4326 (WGS84) para la app",
    "339 poligonos (barrios) en 22 comunas",
    "Columnas: id_barrio, barrio, comuna, shape_leng, shape_area",
    "Area total: ~120 km2",
    "Extension WGS84: lon [-76.5908, -76.4613], lat [3.3318, 3.5019]",
    "Sin IDs duplicados ni nombres vacios",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Estadisticas por comuna", ACCENT_GREEN)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(3.0), [
    "Comuna 19: mas barrios (33, 9.7%)",
    "Comuna 01: menos barrios (4, 1.2%)",
    "Comuna 17: mayor area (12.68 km2, 10.6%)",
    "",
    "Barrio mas grande: Parcelaciones Pance (7.97 km2)",
    "Barrio mas pequeno: La Playa (0.025 km2)",
    "",
    "Visualizaciones:",
    "  Area por comuna (barras)",
    "  Distribucion de areas (histograma)",
    "  Mapa barrios coloreado por comuna",
    "  Mapa comunas disueltas con centroides",
    "  Coropleta de area por barrio",
], sz=11, cc=LIGHT_GRAY)
section(s, Inches(0.8), Inches(4.8), Inches(11), "Integracion con la aplicacion")
bullets(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), [
    "GeoJSON en app/geo-vision-clip-application/public/geo/",
    "Hook useBarriosGeo.ts para carga en frontend",
    "Tooltip hover: nombre de barrio + comuna",
    "Base para analisis de equidad espacial (+4 pts bonif.)",
], sz=11, cc=LIGHT_GRAY)

# === SLIDE 22: HALLAZGOS CLAVE ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 20, "Hallazgos Clave del EDA", "8 conclusiones principales que guian el modelo")
findings = [
    ("Nubosidad critica", "Solo 32.9% escenas S2 y 7-25% pixeles S5P validos.", ACCENT_RED),
    ("Hotspot de aerosoles", "Yumbo-Acopi: AOD ~0.12 superior al promedio.", ACCENT_ORANGE),
    ("BLH: predictor clave", "r = -0.52 con NO2. Covariable N 1.", ACCENT_GREEN),
    ("Separabilidad espectral S2", "NIR y SWIR discriminan 4 clases.", ACCENT_BLUE),
    ("Correlacion NO2-O3", "r = -0.42 (titulacion fotoquimica).", ACCENT_PURPLE),
    ("Tendencia COVID-19", "Caida 18% NO2 en 2020, recuperacion al 92%.", ACCENT_CYAN),
    ("NO2 limitado en DAGMA", "Solo Univalle. Datos = predicciones de modelo.", ACCENT_RED),
    ("Estacionalidad marcada", "NO2 picos jul-ago, O3 mitad de ano, PM secos.", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(findings):
    col, row = i % 2, i // 2
    x = Inches(0.5) + col * Inches(6.3); y = Inches(1.3) + row * Inches(1.45)
    card(s, x, y, Inches(6.0), Inches(1.25))
    rect(s, x, y, Pt(5), Inches(1.25), color)
    tb(s, x+Inches(0.15), y+Inches(0.05), Inches(5.6), Inches(0.25), title, sz=13, cc=color, b=True)
    tb(s, x+Inches(0.15), y+Inches(0.3), Inches(5.6), Inches(0.85), desc, sz=10, cc=LIGHT_GRAY)

# === SLIDE 23: CORRELACIONES ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 21, "Correlaciones Cruzadas", "Entre contaminantes, meteorologia y covariables")
h = ["Variable 1", "Variable 2", "r de Pearson", "Interpretacion"]
r = [
    ["BLH", "NO2", "-0.52", "Mas predictiva. BLH baja atrapa contaminantes"],
    ["T2m", "NO2", "-0.38", "Calor acelera consumo fotoquimico de NO2"],
    ["HR", "O3", "+0.41", "Humedad alta correlaciona con mas ozono"],
    ["Viento", "SO2", "-0.29", "Viento dispersa SO2 industrial"],
    ["NO2", "O3", "-0.42", "Titulacion fotoquimica: NO + O3  NO2 + O2"],
    ["NO2", "SO2", "+0.38", "Fuentes compartidas: transporte + industria"],
    ["T2m", "BLH", "+0.80", "Temperatura y BLH fuertemente acopladas"],
    ["O3", "Temperatura", "+0.60", "Formacion fotoquimica con calor"],
]
tbl(s, Inches(0.4), Inches(1.4), [Inches(2.0), Inches(2.0), Inches(1.5), Inches(5.5)], h, r)
section(s, Inches(0.8), Inches(4.5), Inches(11), "Implicaciones para el modelo")
bullets(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(2.0), [
    "BLH como covariable forzada en ConvLSTM (r=-0.52 justifica su inclusion)",
    "ERA5 no es opcional: la meteorologia modula la contaminacion",
    "Anti-correlacion NO2-O3: el modelo debe aprender dinamica fotoquimica",
    "Correlacion entre contaminantes valida el enfoque multi-output",
    "Separabilidad S2 + correlaciones meteorologicas = base del downscaling",
], sz=12, cc=LIGHT_GRAY)

# === SLIDE 24: GROUND TRUTH LIMITACIONES ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 22, "Ground Truth: Limitaciones y Transparencia", "Datos SVCASC: lo que son y lo que no son")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Que son los datos SVCASC?")
multi(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.0), [
    "SVCASC.FT.50 = Sistema de Vigilancia de Calidad",
    "  del Aire de Santiago de Cali",
    "",
    "NO son mediciones directas de sensores",
    "Son SALIDAS DE MODELOS DE PREDICCION",
    "  Reflejan patrones medios, no eventos agudos",
    "",
    "Ventana limitada: 2020-2023 (4 anos)",
    "  Satelites cubren 2019-2023 (5 anos)",
    "  Extrapolar 2019 = datos sinteticos",
    "  Decision correcta: trabajar con 4 anos reales",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Limitaciones documentadas", ACCENT_RED)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(3.0), [
    "1. NO2 restringido a Univalle",
    "   Imposible validar distribucion espacial",
    "   Solo 1 punto para LOO-CV espacial",
    "",
    "2. Datos generados por modelos",
    "   No capturan episodios agudos reales",
    "   Suavizan la variabilidad real",
    "",
    "3. Desbalance temporal: 4 vs 5 anos",
    "   Reduce ventana de validacion cruzada",
    "",
    "4. Variables con cobertura <5%",
    "   Black Carbon (3%), UV-PM (3%)",
    "   Excluidas del analisis",
], sz=11, cc=LIGHT_GRAY)
section(s, Inches(0.8), Inches(5.0), Inches(11), "Estrategia de mitigacion", ACCENT_ORANGE)
bullets(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.5), [
    "CLIP se entrena con pseudo-labels de S5P (percentiles), NO con DAGMA",
    "DAGMA se usa solo para validacion externa en LOO-CV (Sit. 3)",
    "SVCASC sigue siendo la mejor aproximacion disponible",
], sz=12, cc=LIGHT_GRAY)

# === SLIDE 25: DATASET VISION GENERAL ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 23, "Dataset Situacion 2: Vision General", "2,276 pares imagen-texto  5 clases  Split 70/15/15")
kpi(s, Inches(0.5), Inches(1.4), Inches(2.3), Inches(1.0), "Pares generados", "2,276", ACCENT_GREEN)
kpi(s, Inches(3.0), Inches(1.4), Inches(2.3), Inches(1.0), "Del minimo (1000)", "227%", ACCENT_BLUE)
kpi(s, Inches(5.5), Inches(1.4), Inches(2.3), Inches(1.0), "Clases", "5", ACCENT_ORANGE)
kpi(s, Inches(8.0), Inches(1.4), Inches(2.3), Inches(1.0), "Secuencias 8 fechas", "30", ACCENT_PURPLE)
kpi(s, Inches(10.5), Inches(1.4), Inches(2.3), Inches(1.0), "Split 70/15/15", "SEED=42", ACCENT_CYAN)
section(s, Inches(0.8), Inches(2.7), Inches(6), "Conceptos clave")
bullets(s, Inches(0.8), Inches(3.1), Inches(5.5), Inches(3.0), [
    "Tile: recorte 64x64 px (~640 m) de escena Sentinel-2",
    "13 bandas (B1-B12 + SCL), resolucion 10/20/60 m",
    "1 escena S2  hasta 40 tiles (stride 32 px)",
    "Descripcion: texto en espanol por plantillas v2",
    "Pseudo-label: percentil S5P sobre el centroide del tile",
    "Ventana S5P: +/-7 dias alrededor de fecha S2 (nanmax)",
], sz=11)
section(s, Inches(7), Inches(2.7), Inches(5.5), "Las 5 clases", ACCENT_GREEN)
h = ["Clase", "Cantidad", "Criterio"]
r = [
    ["contaminacion_alta_NO2", "500", "Exceso NO2 vs percentil Cali"],
    ["contaminacion_alta_SO2", "276", "Score v2: mayor exceso relativo"],
    ["ozono_anomalo", "500", "O3 vs p90/p25 segun contexto"],
    ["vegetacion_densa", "500", "NDVI >= 0.45, BSI bajo"],
    ["suelo_urbano", "500", "NDVI <= 0.35, NDBI alto"],
]
tbl(s, Inches(7), Inches(3.1), [Inches(2.8), Inches(1.2), Inches(2.5)], h, r)

# === SLIDE 26: PIPELINE GENERACION ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 24, "Pipeline de Generacion del Dataset", "generar_dataset_sit2_v2.py  1,218 lineas  Score-based")
steps = [
    ("1. Percentiles", "Lee Zarr S5P de GCS\nStep=8 temporal, step=2\np25/p50/p75/p90/p99", ACCENT_BLUE),
    ("2. Barrido S2", "Itera escenas S2\nThreadPoolExecutor\nFiltro SCL nube/claro", ACCENT_GREEN),
    ("3. Tiles 64x64", "Genera tiles con stride\nEvalua cada candidato\nNDVI, BSI, NDBI", ACCENT_ORANGE),
    ("4. Clase score", "Score = valor/umbral p90\nMaximo score  clase\nBalance tracking", ACCENT_RED),
    ("5. Descripcion", "Plantillas v2 distintivas\nIncluye exceso_*, ndbi\nFecha, lat, lon", ACCENT_PURPLE),
    ("6. Split+flush", "70/15/15 estratificado\nSecuencias temporales\nZarr incremental", ACCENT_CYAN),
]
for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.2) + i * Inches(2.2)
    card(s, x, Inches(1.7), Inches(2.0), Inches(1.5))
    rect(s, x, Inches(1.7), Inches(2.0), Pt(4), color)
    tb(s, x+Inches(0.08), Inches(1.78), Inches(1.85), Inches(0.25), title, sz=10, cc=color, b=True, al=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.08), Inches(2.1), Inches(1.85), Inches(0.95), desc, sz=8, cc=LIGHT_GRAY, al=PP_ALIGN.CENTER)
    if i < 5:
        tb(s, x+Inches(2.0), Inches(2.2), Inches(0.2), Inches(0.4), ">", sz=16, cc=MEDIUM_GRAY, b=True)
section(s, Inches(0.8), Inches(3.5), Inches(11), "Clase _AccesoS5P  Caching inteligente")
multi(s, Inches(0.8), Inches(3.9), Inches(5.5), Inches(3.0), [
    "Cachea los 3 Zarr de S5P (NO2, SO2, O3) en memoria",
    "mapa_compuesto(fecha)  mapa 2D por contaminante",
    "valor_en(lat, lon, fecha)  punto exacto o mediana",
    "  buffer 2px si no cae exacto en la grilla",
    "nanmax para agregar multiples orbitas en +/-7 dias",
], sz=11)
section(s, Inches(7), Inches(3.5), Inches(5.5), "Clase IncrementalTilesZarr", ACCENT_GREEN)
multi(s, Inches(7), Inches(3.9), Inches(5.5), Inches(3.0), [
    "Escribe tiles.zarr por lotes, sin acumular todo en RAM",
    "Buffer de N tiles + flush periodico",
    "Redimensiona el Zarr conforme crece (resize)",
    "Thread-safe para escritura concurrente",
], sz=11, cc=LIGHT_GRAY)

# === SLIDE 27: v1 vs v2 ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 25, "Dataset: v1 vs v2", "Evolucion del pipeline de generacion")
h = ["Caracteristica", "v1 (par_imagen_texto.py)", "v2 (dataset_sit2_v2.py)"]
r = [
    ["Asignacion clase", "Prioridad fija NO2>SO2>O3>Veg>Urb", "Score-based: max exceso relativo"],
    ["Indices espectrales", "NDVI, BSI", "NDVI, BSI, NDBI"],
    ["Columnas diagnostico", "No", "exceso_*, ndbi, p90_*"],
    ["Descripciones", "Genericas, poco distintivas", "Mas distintivas, incluyen scores"],
    ["Balance tracking", "No (limite duro por clase)", "Pre-cap por escena + tracking"],
    ["Lineas codigo", "1,406", "1,218"],
]
tbl(s, Inches(0.4), Inches(1.5), [Inches(2.5), Inches(4.5), Inches(4.5)], h, r)
section(s, Inches(0.8), Inches(4.2), Inches(11), "Impacto de las mejoras v2")
bullets(s, Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.0), [
    "Score-based elimina sesgo de prioridad fija (NO2 ya no siempre gana)",
    "NDBI mejora discriminacion urbano vs industrial vs vegetacion",
    "Descripciones mas distintivas  embeddings de texto mas separables",
    "Columnas de diagnostico permiten analizar causas de error",
], sz=12, cc=LIGHT_GRAY)

# === SLIDE 28: PERCENTILES ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 26, "Percentiles Globales S5P (Cali, 5 anos)", "Umbrales para etiquetado semi-supervisado")
h = ["Contaminante", "p25", "p50", "p75", "p90", "p99"]
r = [
    ["NO2 (mol/m2)", "1.76e-5", "2.40e-5", "3.27e-5", "4.57e-5", "8.56e-5"],
    ["SO2 (mol/m2)", "1.27e-4", "2.04e-4", "3.25e-4", "5.90e-4", "1.28e-3"],
    ["O3 (mol/m2)", "0.113", "0.115", "0.118", "0.122", "0.129"],
]
tbl(s, Inches(1.5), Inches(1.5), [Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)], h, r)
section(s, Inches(0.8), Inches(3.0), Inches(11), "Como se usan los percentiles")
bullets(s, Inches(0.8), Inches(3.4), Inches(11.5), Inches(3.5), [
    "Muestreo: step temporal 8, step espacial 2 sobre Zarr",
    "Ventana temporal: +/-7 dias alrededor de fecha S2",
    "Agregacion: nanmax (maximo de orbitas disponibles)",
    "",
    "Regla de etiquetado v2 (score-based):",
    "  score_NO2 = valor_NO2 / p90_NO2",
    "  score_SO2 = valor_SO2 / p90_SO2",
    "  score_O3 = valor_O3 / p90_O3",
    "  Clase = argmax(score) si max(score) > 1.0",
    "  Si no: evaluar NDVI/BSI/NDBI para vegetacion/urbano",
    "",
    "Ventaja: NO2 deja de tener prioridad absoluta sobre SO2/O3",
], sz=12, cc=LIGHT_GRAY)

# === SLIDE 29: ESTRUCTURA DATASET ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 27, "Estructura del Dataset Generado", "Archivos de salida en dataset_sit2/")
section(s, Inches(0.8), Inches(1.3), Inches(5.5), "Archivos")
bullets(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.5), [
    "tiles.zarr: (2276, 13, 64, 64) int16",
    "metadatos.parquet: 2276 filas",
    "percentiles.json: umbrales p25-p99",
    "secuencias.json: 30 secuencias x 8 fechas",
    "resumen.json: estadisticas globales",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Columnas del Parquet", ACCENT_GREEN)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(3.0), [
    "tile_id, clase, descripcion, split",
    "centroide_lat, centroide_lon",
    "ndvi, bsi, ndbi",
    "no2, so2, o3 (valores S5P)",
    "exceso_NO2, exceso_SO2, exceso_O3",
    "p90_NO2, p90_SO2, p90_O3",
    "fecha, frac_nubes, frac_claros",
    "valid_ratio",
], sz=10, cc=LIGHT_GRAY)
section(s, Inches(0.8), Inches(4.5), Inches(11), "Split estratificado (70/15/15)")
bullets(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(2.0), [
    "Train: 1,593 tiles  Val: 341 tiles  Test: 342 tiles",
    "Estratificado por clase: cada clase mantiene proporcion",
    "Semilla fija SEED=42  reproducible entre integrantes",
    "Secuencias temporales agrupadas por celda 0.1 grado",
    "30 secuencias de 8 fechas cuasi-consecutivas para ConvLSTM",
], sz=12, cc=LIGHT_GRAY)

# === SLIDE 30: EDA DATASET ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 28, "EDA del Dataset: Calidad y Problemas", "sit2_EDA_pares_dataset.ipynb  2,613 lineas")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Analisis de calidad")
bullets(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.5), [
    "Balance global y por split (barras por clase)",
    "Estadisticas por clase: NDVI, BSI, NO2, SO2 (boxplots)",
    "Matriz de correlacion entre variables numericas",
    "Distribucion temporal de adquisiciones por split",
    "Muestras RGB: 4 tiles aleatorios por clase",
    "146 tiles de baja calidad identificados",
    "Filtro estricto: 88.91% pasa (nubes<=0.15, claros>=0.90)",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Problemas detectados", ACCENT_RED)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(3.0), [
    "DATA LEAKAGE: 100% escenas val en train",
    "  180 escenas S2 compartidas train/val",
    "  Recall optimista: el modelo memoriza",
    "",
    "Textos casi identicos dentro de cada clase",
    "  Embeddings de texto colapsan",
    "",
    "~38 tiles con clase inconsistente con lo visual",
    "  Ruido en pseudo-labels por escala S5P",
    "",
    "Ozono alto vs bajo: mismo template textual",
    "",
    "~4.5% tiles casi negros (poca informacion)",
], sz=10, cc=LIGHT_GRAY)
section(s, Inches(0.8), Inches(5.0), Inches(11), "Completitud temporal para Sit. 3")
bullets(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.5), [
    "Grilla 0.05 grados  ~863 secuencias deslizantes posibles",
    "30 secuencias de 8 fechas generadas (minimo consigna)",
], sz=11, cc=LIGHT_GRAY)

# === SLIDE 31: DIAGNOSTICO Y MEJORAS ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 29, "Diagnostico y Plan de Mejora", "INFORME_DIAGNOSTICO_SIT2.md  340 lineas")
causas = [
    ("Causa 1", "Textos no discriminativos", "Embeddings de texto de la misma clase casi identicos  espacio contrastivo no puede alinear imagenes distintas con textos indistinguibles.", ACCENT_RED),
    ("Causa 2", "Data leakage train/val", "100% escenas val en train  modelo memoriza en lugar de generalizar  Recall optimista.", ACCENT_ORANGE),
    ("Causa 3", "SAE lineal sin ReLU", "Encoder lineal produce z gaussiano  sparsity techo ~5%. Sparsity 70% imposible sin no-linealidad.", ACCENT_PURPLE),
]
for i, (tit, subtit, desc, color) in enumerate(causas):
    x = Inches(0.5) + i * Inches(4.2)
    card(s, x, Inches(1.7), Inches(3.95), Inches(1.8))
    rect(s, x, Inches(1.7), Inches(3.95), Pt(4), color)
    tb(s, x+Inches(0.12), Inches(1.78), Inches(3.7), Inches(0.25), tit, sz=11, cc=color, b=True)
    tb(s, x+Inches(0.12), Inches(2.0), Inches(3.7), Inches(0.2), subtit, sz=10, cc=WHITE, b=True)
    tb(s, x+Inches(0.12), Inches(2.25), Inches(3.7), Inches(1.1), desc, sz=9, cc=LIGHT_GRAY)

section(s, Inches(0.8), Inches(3.8), Inches(11), "Recomendaciones del informe diagnostico")
recoms = [
    ("R-D1", "Split por escena S2, no por tile  elimina leakage", ACCENT_GREEN),
    ("R-D2", "Filtros SCL mas estrictos: solo pixeles claros (clases 4-7)", ACCENT_GREEN),
    ("R-D3", "Aumentar diversidad de descripciones (mas plantillas)", ACCENT_GREEN),
    ("R-D4", "Verificar coherencia visual-S5P con buffer espacial mayor", ACCENT_GREEN),
    ("R-D5", "Expandir ventana S5P a +/-14 dias para mejorar cobertura", ACCENT_GREEN),
    ("R-M1", "ReLU en encoder SAE + lambda L1 mas agresivo (sparsity >=70%)", ACCENT_BLUE),
    ("R-M2", "Split por escena + early stopping con paciencia reducida", ACCENT_BLUE),
    ("R-M3", "LoRA rank 16 en bloques 6-11 de ViT y texto", ACCENT_BLUE),
    ("R-M4", "12 bandas opticas con conv1 multiespectral (no true-color)", ACCENT_BLUE),
    ("R-M5", "Prompts ES+EN: ~10 variantes por clase para diversidad textual", ACCENT_BLUE),
]
for i, (code, rec, color) in enumerate(recoms):
    col, row = i % 2, i // 2
    x = Inches(0.5) + col * Inches(6.3); y = Inches(4.2) + row * Inches(0.42)
    tb(s, x, y, Inches(0.5), Inches(0.35), code, sz=9, cc=color, b=True)
    tb(s, x+Inches(0.5), y, Inches(5.5), Inches(0.35), rec, sz=10, cc=LIGHT_GRAY)

# === SLIDE 32: MODELO ARQUITECTURA ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 30, "Modelo GeoVision-CLIP + SAE", "Arquitectura multimodal contrastiva")
blocks = [
    ("Imagen S2\nTile 64x64\n13 bandas", ACCENT_BLUE),
    ("Conv 13->3\ntrue-color\nB4,B3,B2", ACCENT_CYAN),
    ("ViT-B/32\nRemoteCLIP\n768-d frozen", ACCENT_GREEN),
    ("SAE_img\n512->512\nSparse AE", ACCENT_ORANGE),
    ("Proj Head\n512->256\nEspacio C", ACCENT_RED),
]
for i, (txt, color) in enumerate(blocks):
    x = Inches(0.3) + i * Inches(2.55)
    card(s, x, Inches(1.5), Inches(2.3), Inches(1.3))
    rect(s, x, Inches(1.5), Inches(2.3), Pt(4), color)
    tb(s, x+Inches(0.05), Inches(1.6), Inches(2.2), Inches(1.0), txt, sz=10, cc=WHITE, al=PP_ALIGN.CENTER)
    if i < 4:
        tb(s, x+Inches(2.3), Inches(1.9), Inches(0.3), Inches(0.4), ">", sz=18, cc=MEDIUM_GRAY, b=True)
blocks_t = [
    ("Texto ES\nDescripcion\nmax 256 tok", ACCENT_GREEN),
    ("MiniLM-L12\n384-d\nMultilingue", ACCENT_CYAN),
    ("Linear\n384->512", ACCENT_ORANGE),
    ("SAE_txt\n512->512\nSparse AE", ACCENT_ORANGE),
    ("Proj Head\n512->256\nEspacio C", ACCENT_RED),
]
for i, (txt, color) in enumerate(blocks_t):
    x = Inches(0.3) + i * Inches(2.55); y = Inches(3.2)
    card(s, x, y, Inches(2.3), Inches(1.3))
    rect(s, x, y, Inches(2.3), Pt(4), color)
    tb(s, x+Inches(0.05), y+Inches(0.1), Inches(2.2), Inches(1.0), txt, sz=10, cc=WHITE, al=PP_ALIGN.CENTER)
    if i < 4:
        tb(s, x+Inches(2.3), y+Inches(0.4), Inches(0.3), Inches(0.4), ">", sz=18, cc=MEDIUM_GRAY, b=True)
card(s, Inches(4.5), Inches(4.8), Inches(4.3), Inches(1.2))
rect(s, Inches(4.5), Inches(4.8), Inches(4.3), Pt(4), ACCENT_GREEN)
tb(s, Inches(4.6), Inches(4.9), Inches(4.1), Inches(0.3), "InfoNCE Contrastive Loss", sz=12, cc=ACCENT_GREEN, b=True, al=PP_ALIGN.CENTER)
tb(s, Inches(4.6), Inches(5.2), Inches(4.1), Inches(0.7),
   "L = L_InfoNCE + a*(L_SAE_img + L_SAE_txt)\nL_SAE = MSE(x, xh) + l*||z||1\nt=0.07 (learnable)  a=0.1  l=1e-3\nd=512 SAE  dim contrastiva=256",
   sz=9, cc=LIGHT_GRAY, al=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.4),
   "Frozen parcial: ViT congelado, fine-tune condicional en capas finales. SAE y projection heads entrenables.", sz=11, cc=ACCENT_ORANGE)

# === SLIDE 33: RAMA IMAGEN ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 31, "Rama Imagen: RemoteCLIP + SAE", "Encoder visual con adaptacion multiespectral")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Arquitectura detallada")
multi(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.0), [
    "Entrada: Tile (13, 64, 64)  13 bandas Sentinel-2",
    "",
    "Paso 1  Adaptacion espectral:",
    "  Conv2d 3->12 canales (pesos RemoteCLIP)",
    "  True-color: B4->R, B3->G, B2->B",
    "  9 canales restantes: media RGB * (3/12)",
    "  Interpolacion bilineal a 224x224 px",
    "",
    "Paso 2  ViT-B/32 RemoteCLIP:",
    "  Congelado (frozen) durante entrenamiento",
    "  Fine-tune condicional con LoRA",
    "  Salida: embedding 768-d",
    "",
    "Paso 3  SAE visual:",
    "  Encoder: Linear(768, 512)",
    "  Decoder: Linear(512, 768)",
    "  Sparsity target: >70%",
    "",
    "Paso 4  Projection head:",
    "  Linear(512, 256) sin bias",
    "  Normalizacion L2  espacio unitario",
], sz=10, cc=LIGHT_GRAY)
section(s, Inches(7), Inches(1.3), Inches(5.5), "LoRA adapters", ACCENT_ORANGE)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(2.5), [
    "Implementacion: LoRALinear (nn.Module custom)",
    "Rank: 8 (experimental)  16 (en mejora)",
    "Aplica desde bloque 6 del ViT en adelante",
    "Solo pesos query/value de atencion",
    "Inferencia: Wx + (B*A)x, sin overhead extra",
    "",
    "Ventaja: adaptar RemoteCLIP al dominio",
    "  Sentinel-2 sin perder representaciones",
    "  pre-entrenadas en imagenes de satelite",
], sz=10, cc=LIGHT_GRAY)

# === SLIDE 34: RAMA TEXTO ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 32, "Rama Texto: MiniLM + SAE", "Encoder textual multilingue para descripciones en espanol")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Arquitectura detallada")
multi(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.0), [
    "Entrada: Descripcion textual en espanol",
    "  Max tokens: 256",
    "  Tokenizer: multilingue (sentencepiece)",
    "",
    "Paso 1  MiniLM-L12:",
    "  paraphrase-multilingual-MiniLM",
    "  12 capas, 384-d hidden",
    "  Congelado (frozen): solo LoRA",
    "  Salida: embedding 384-d",
    "",
    "Paso 2  Linear adapter:",
    "  Linear(384, 512) para igualar SAE",
    "",
    "Paso 3  SAE texto:",
    "  Encoder: Linear(512, 512)",
    "  Decoder: Linear(512, 512)",
    "  Misma arquitectura que SAE visual",
    "",
    "Paso 4  Projection head:",
    "  Linear(512, 256) sin bias",
    "  Normalizacion L2  mismo espacio",
], sz=10, cc=LIGHT_GRAY)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Ejemplos descripciones v2", ACCENT_GREEN)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(4.5), [
    "contaminacion_alta_NO2:",
    "  'Imagen satelital de Cali con contaminacion",
    "  alta de NO2. Nivel excede el percentil 90",
    "  de la distribucion historica (4.57e-5).'",
    "",
    "contaminacion_alta_SO2:",
    "  'Tile con concentracion critica de SO2.",
    "  Valor: 6.02e-4 mol/m2. Probable fuente:",
    "  corredor industrial Yumbo-Acopi.'",
    "",
    "vegetacion_densa:",
    "  'Cobertura vegetal densa con NDVI de 0.72.",
    "  BSI bajo (-0.15). Zona rural o parque.'",
    "",
    "suelo_urbano:",
    "  'Zona urbana consolidada. NDBI: 0.18.",
    "  Fraccion construida alta. Bajo NDVI.'",
], sz=9, cc=LIGHT_GRAY)

# === SLIDE 35: SAE POST-TRAIN ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 33, "Sparse Autoencoder (Post-train)", "Entrenamiento separado sobre embeddings congelados")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Arquitectura y parametros")
multi(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.5), [
    "Arquitectura: 512  ReLU  2048  Linear  512",
    "  d_hidden = 4 x d_model = 2048",
    "",
    "Hiperparametros:",
    "  EPOCHS = 8,000 (ocho mil)",
    "  LR = 1e-3, AdamW, CosineAnnealingLR",
    "  LAMBDA_L1 = 1e-2 (sparsity)",
    "  BATCH_SIZE = 256",
    "  CLIP_GRAD_NORM = 1.0",
    "",
    "Decoder normalization: filas a norma unitaria",
    "Adaptive lambda:",
    "  sparsity < target-0.05  l x1.15",
    "  sparsity > target+0.1   l x0.85",
], sz=10, cc=LIGHT_GRAY)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Resultados espectaculares", ACCENT_GREEN)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(2.0), [
    "MSE minimo: 1.29e-06 (KPI <= 0.05  EXCELENTE)",
    "Sparsity final: 94.1% (KPI >= 0.70  EXCELENTE)",
    "",
    "Evolucion del entrenamiento:",
    "  Epoca 0:    MSE 0.0014, sparsity 63.7%",
    "  Epoca 800:  MSE 8e-6,   sparsity 92.9%",
    "  Epoca 7999: MSE 1.3e-6, sparsity 94.1%",
], sz=10, cc=LIGHT_GRAY)
section(s, Inches(0.8), Inches(4.0), Inches(11), "Analisis de neuronas (interpretabilidad mecanica)")
h = ["Neurona", "Correlacion mas alta", "r", "Interpretacion"]
r = [
    ["1415", "NDVI", "+0.8192", "Detector de vegetacion"],
    ["1907", "NDVI", "-0.7171", "Detector de ausencia vegetacion"],
    ["1415", "BSI", "-0.7912", "Anti-suelo desnudo"],
    ["1907", "BSI", "+0.7504", "Detector de suelo desnudo"],
    ["116", "NO2", "+0.3451", "Detector de contaminacion"],
    ["1226", "frac_nodata", "+0.4971", "Detector de datos faltantes"],
]
tbl(s, Inches(0.3), Inches(4.3), [Inches(1.2), Inches(2.5), Inches(1.0), Inches(3.5)], h, r)

# === SLIDE 36: ENTRENAMIENTO CLIP ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 34, "Entrenamiento CLIP (Sit. 2)", "20 epochs  LoRA  RemoteCLIP adaptado")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Configuracion")
multi(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.0), [
    "Framework: PyTorch Lightning",
    "GPU: T4/P100 (Kaggle), A100 (Colab)",
    "Batch size: 32-64",
    "Epochs: 20 (early stopping epoca 8)",
    "Optimizer: AdamW (lr=1e-5, wd=0.2)",
    "Scheduler: CosineAnnealingLR",
    "Loss: CrossEntropy + label_smoothing=0.1",
    "LoRA: rank 8, bloques 6-11 del ViT",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Curvas de entrenamiento", ACCENT_GREEN)
h = ["Epoca", "Train CE", "Val CE", "Acc", "R@1", "R@5"]
r = [
    ["0", "1.5483", "1.4804", "0.363", "0.363", "1.0"],
    ["1", "1.4267", "1.4457", "0.407", "0.395", "1.0"],
    ["2", "1.3899", "1.4240", "0.434", "0.440", "1.0"],
    ["8", "1.0441", "1.4705", "0.434", "0.434", "1.0"],
    ["19", "0.5114", "1.7130", "0.386", "0.378", "1.0"],
]
tbl(s, Inches(7), Inches(1.7), [Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.0), Inches(1.0), Inches(1.0)], h, r)
section(s, Inches(0.8), Inches(3.5), Inches(11), "Analisis de resultados")
bullets(s, Inches(0.8), Inches(3.9), Inches(11.5), Inches(3.0), [
    "Mejor accuracy val: 43.4% (epoca 8)  luego overfitting",
    "Test accuracy: 47.4%, Test Recall@1: 46.8%, Test R@5: 100%",
    "R@5=100% es trivial (solo 5 clases)",
    "R@1=0.47 cerca del umbral 0.45, pero hay data leakage",
    "Overfitting: train loss 0.51, val loss sube a 1.71",
    "Matriz confusion: vegetacion 52%, NO2 44%, SO2 28%",
    "SO2 confunde con NO2 (misma fuente industrial)",
], sz=11, cc=LIGHT_GRAY)

# === SLIDE 37: KPIS ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 35, "KPIs: Estado vs Metas", "Comparacion consigna vs logrado con semaforo")
h = ["KPI", "Minimo", "Excelente", "Logrado", "Estado"]
r = [
    ["Recall@1 imagen->texto", ">= 0.45", ">= 0.65", "0.11 (CLIP) / 0.47 (CNN)", "En mejora"],
    ["Recall@5 imagen->texto", ">= 0.70", ">= 0.85", "0.24 (CLIP) / 1.0 (CNN)", "En mejora"],
    ["Sparsity ratio SAE visual", ">= 0.70", ">= 0.85", "0.941 (post-train)", "EXCELENTE"],
    ["Loss reconstruccion SAE", "<= 0.05", "<= 0.02", "1.29e-06", "EXCELENTE"],
    ["Varianza explicada AFE", ">= 80%", ">= 90%", "85.56% (2 factores)", "LOGRADO"],
    ["RMSEA (modelo AFC)", "< 0.08", "< 0.05", "0.0417", "EXCELENTE"],
    ["CFI (modelo AFC)", "> 0.90", "> 0.95", "0.9989", "EXCELENTE"],
    ["Pares imagen-texto", ">= 1,000", ">= 1,500", "2,276", "EXCELENTE"],
    ["Secuencias >=8 fechas", ">= 30", ">= 50", "30", "LOGRADO"],
    ["Panel >=50 GB", ">= 50 GB", ">= 60 GB", "66.78 GB", "EXCELENTE"],
]
tbl(s, Inches(0.3), Inches(1.4), [Inches(3.2), Inches(1.5), Inches(1.5), Inches(3.0), Inches(2.0)], h, r)
section(s, Inches(0.8), Inches(5.0), Inches(11), "Notas sobre los KPIs")
bullets(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.5), [
    "Recall@1 CLIP puro (0.11) vs CNN classifier (0.47): CLIP contrastivo en etapa temprana",
    "R@5=1.0 del CNN es trivial (solo 5 clases)  la metrica relevante es R@1",
    "SAE post-train independiente ya cumple TODOS los KPIs (sparsity, MSE, AFE, AFC)",
    "AFE+AFC realizados sobre embeddings SAE, no sobre CLIP raw",
], sz=12, cc=LIGHT_GRAY)

# === SLIDE 38: AFE + AFC ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 36, "Validacion Psicometrica: AFE + AFC", "Analisis factorial sobre embeddings SAE")
section(s, Inches(0.8), Inches(1.3), Inches(5.5), "AFE (Analisis Factorial Exploratorio)")
multi(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.5), [
    "Metodo: FactorAnalyzer con rotacion Varimax",
    "Extraccion: Componentes Principales",
    "",
    "Resultados:",
    "  KMO = 0.9746  EXCELENTE (>0.9 = marvelous)",
    "  Bartlett chi2 = 103,790, p = 0.0",
    "  Criterio Kaiser: 2 factores retenidos",
    "  Varianza explicada: 85.56%  KPI>=80% OK",
    "  RMSR = 0.0206  EXCELENTE",
    "  Scree plot con caida clara tras factor 2",
    "",
    "Interpretacion de los 2 factores:",
    "  Factor 1: Carga Urbana / Antropogenica",
    "  Factor 2: Estres Vegetal / Atmosferico",
], sz=10, cc=LIGHT_GRAY)
section(s, Inches(7), Inches(1.3), Inches(5.5), "AFC (Analisis Factorial Confirmatorio)", ACCENT_GREEN)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(3.5), [
    "Software: semopy (SEM en Python)",
    "Estimador: DWLS (robusto, datos no normales)",
    "",
    "Modelo:",
    "  Carga_Urbana  7 indicadores",
    "  Estres_Vegetal_Atmosferico  7 indicadores",
    "",
    "Resultados:",
    "  RMSEA = 0.0417  KPI <0.08 OK  EXCELENTE",
    "  CFI = 0.9989    KPI >0.90 OK  EXCELENTE",
    "  chi2 = 375.41, gl = 76",
    "",
    "Los embeddings SAE tienen estructura",
    "factorial consistente con constructos",
    "teoricos: el espacio latente es",
    "psicometricamente valido",
], sz=10, cc=LIGHT_GRAY)
card(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.0))
tb(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(0.8),
   "Los SAE aprenden representaciones que (1) reconstruyen con MSE 1e-6, (2) son 94% dispersas, (3) se organizan en 2 factores interpretables (85.6% varianza), (4) pasan pruebas estrictas de validez de constructo (RMSEA=0.04, CFI=0.999).",
   sz=11, cc=LIGHT_GRAY)

# === SLIDE 39: SIT3 TENSOR ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 37, "Situacion 3: Generacion Tensor ConvLSTM", "sit3_generar_tensor_convlstm.ipynb")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Pipeline de generacion")
bullets(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.0), [
    "1. Descarga dataset_sit2 desde HuggingFace Hub",
    "2. Carga S5P (NO2, SO2, O3) desde Zarr en GCS",
    "3. Carga checkpoint RemoteCLIP fine-tuned (best.pt)",
    "4. Adaptacion conv1: 3->12 canales",
    "5. Extraccion embeddings 512-d:",
    "   Normalizacion sample-wise (mean/std 512 tiles)",
    "   Batch processing (B=64)",
    "   Interpolacion bilineal a 224x224",
    "6. Generacion secuencias: agrupa por celda 0.05 grados",
    "7. Targets S5P: T+1, T+3, T+7 (nearest, tol 3 dias)",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Arquitectura del tensor", ACCENT_GREEN)
multi(s, Inches(7), Inches(1.7), Inches(5.5), Inches(3.5), [
    "Entrada X: (N, 8, 522, 5, 5)",
    "  N = numero de secuencias",
    "  8 = pasos temporales (SEQ_LEN)",
    "  522 canales:",
    "    512 = embeddings RemoteCLIP",
    "    10 = codificacion positional:",
    "      doy_sin, doy_cos, mes_sin, mes_cos",
    "      lat_norm, lon_norm, time_delta",
    "      ndvi, bsi, ndbi",
    "  5x5 = grilla espacial (0.005 deg res)",
    "",
    "Target y: (N, 3, 3)",
    "  3 horizontes: T+1, T+3, T+7",
    "  3 contaminantes: NO2, SO2, O3",
], sz=10, cc=LIGHT_GRAY)
section(s, Inches(0.8), Inches(5.0), Inches(11), "Configuracion")
bullets(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.5), [
    "SEQ_LEN=8, HORIZONS=[1,3,7], GS=5, STRIDE=0.005 deg",
    "C=522, DEVICE=cuda, subido a HF: Slucu-0310/geovision-cali-sit3",
    "Tensor listo para ConvLSTM bidireccional (hidden=128, kernel=3, 2 capas)",
], sz=12, cc=LIGHT_GRAY)

# === SLIDE 40: COVARIABLES ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 38, "Covariables para ConvLSTM", "preparar_covariables_convlstm.py  366 lineas")
h = ["Feature", "Descripcion", "Tipo"]
r = [
    ["doy_sin, doy_cos", "Dia del ano (codif. circular)", "Temporal"],
    ["mes_sin, mes_cos", "Mes del ano (codif. circular)", "Temporal"],
    ["lat_norm, lon_norm", "Coordenadas normalizadas", "Espacial"],
    ["no2_norm, so2_norm, o3_norm", "Normalizados (p50/p90)", "Quimico"],
    ["no2_log, so2_log, o3_log", "Log-transformacion", "Quimico"],
    ["no2_ugm3, so2_ugm3, o3_ugm3", "mol/m2  ug/m3", "Quimico"],
    ["gap_dias_desde_anterior", "Dias desde tile previo", "Temporal"],
    ["era5_t2m", "Temperatura 2m (ERA5)", "Meteo"],
    ["era5_blh", "Boundary layer height", "Meteo"],
    ["era5_wind_u, era5_wind_v", "Viento componentes", "Meteo"],
    ["era5_presion", "Presion superficial", "Meteo"],
    ["modis_aod_047, modis_aod_055", "AOD (MODIS)", "Aerosol"],
]
tbl(s, Inches(0.3), Inches(1.5), [Inches(2.8), Inches(4.5), Inches(2.0)], h, r)
section(s, Inches(0.8), Inches(5.5), Inches(11), "Conversiones y normalizacion")
bullets(s, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.0), [
    "mol/m2  ug/m3: asumiendo 8 km de troposfera para Cali",
    "Normalizacion robusta: (valor - p50) / (p90 - p50)",
    "NaN imputados con p50 del contaminante",
    "ERA5 y MODIS: nearest neighbor desde Zarr GCS (opcional)",
], sz=11, cc=LIGHT_GRAY)

# === SLIDE 41: ARQUITECTURA PIPELINE (imagen) ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 39, "Arquitectura del Pipeline", "Diagrama completo del sistema")
arq_img = os.path.join(BASE_DIR, "notebooks", "Situacion_1", "Arquitectura", "arquitecturaSIT1.png")
if os.path.exists(arq_img):
    s.shapes.add_picture(arq_img, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
else:
    tb(s, Inches(1), Inches(3), Inches(11), Inches(1.5),
       "Diagrama de arquitectura no encontrado", sz=18, cc=ACCENT_RED, al=PP_ALIGN.CENTER)

# === SLIDE 42: FRONTEND ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 40, "Frontend Next.js + Leaflet", "App profesional  No Streamlit  Modo oscuro")
section(s, Inches(0.8), Inches(1.3), Inches(5.5), "Stack del frontend")
multi(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.0), [
    "Next.js 16 + React 19 + TypeScript",
    "Leaflet / React-Leaflet para mapas",
    "Zustand (estado global)",
    "React Query (datos asincronos)",
    "shadcn/ui + Tailwind CSS",
    "Modo oscuro (+2 pts bonificacion)",
    "",
    "Ruta: app/geo-vision-clip-application/",
    "",
    "Hooks personalizados:",
    "  useStations  9 estaciones DAGMA",
    "  usePrediction  consulta /predict",
    "  useBarriosGeo  339 barrios, 22 comunas",
    "  useValidation  vs ground truth",
    "  useScrollReveal  animaciones",
], sz=11, cc=LIGHT_GRAY)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Experiencia de usuario", ACCENT_GREEN)
interactions = [
    ("Click en mapa", "Prediccion en (lat, lon) + tooltip valor +/- s"),
    ("Selector contaminante", "NO2  SO2  O3"),
    ("Slider temporal", "T+1  T+3  T+7"),
    ("Hover poligono", "Tooltip con barrio + comuna"),
    ("Toggle capas", "Limites administrativos ON/OFF"),
    ("Descarga", "Prediccion como GeoTIFF o CSV"),
]
for i, (action, result) in enumerate(interactions):
    y = Inches(1.7) + i * Inches(0.55)
    card(s, Inches(7), y, Inches(5.5), Inches(0.45))
    tb(s, Inches(7.15), y+Inches(0.02), Inches(1.8), Inches(0.25), action, sz=11, cc=ACCENT_BLUE, b=True)
    tb(s, Inches(7.15), y+Inches(0.24), Inches(5.2), Inches(0.2), result, sz=10, cc=LIGHT_GRAY)
section(s, Inches(7), Inches(5.2), Inches(5.5), "Datos geograficos integrados", ACCENT_ORANGE)
bullets(s, Inches(7), Inches(5.6), Inches(5.5), Inches(1.5), [
    "comunas_cali.geojson  22 comunas oficiales",
    "barrios_cali.geojson  339 barrios (WGS84)",
    "manifest.json  metadatos de capas geo",
], sz=11, cc=LIGHT_GRAY)

# === SLIDE 43: PUBLICACION HF ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 41, "Publicacion en HuggingFace Hub", "Dataset disponible para reproducibilidad")
section(s, Inches(0.8), Inches(1.3), Inches(6), "Repositorios publicados")
bullets(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.0), [
    "HF: Slucu-0310/geovision-cali-sit2",
    "  tiles.zarr + metadatos.parquet",
    "  percentiles.json + secuencias.json",
    "  Acceso directo desde Colab/Kaggle",
    "",
    "HF: Slucu-0310/geovision-cali-sit3",
    "  Tensor ConvLSTM (N, 8, 522, 5, 5)",
    "  Listo para entrenamiento",
    "",
    "publicar_hf.py:",
    "  Modo subir: GCS  HF Hub",
    "  Modo descargar: HF Hub  local",
    "  Excluye TIFFs crudos automaticamente",
    "  Modo ligero: sin Sentinel-2 (~66 GB)",
], sz=11)
section(s, Inches(7), Inches(1.3), Inches(5.5), "Artefactos del modelo", ACCENT_GREEN)
bullets(s, Inches(7), Inches(1.7), Inches(5.5), Inches(2.0), [
    "Checkpoint CLIP: best.pt",
    "  SHA256: 2f65a2fd41447f58b1ef8c14daff574d...",
    "  Incluye: metadata.json, hashes.sha256",
    "  training_logs.json, training_history.json",
    "",
    "Checkpoint SAE: sae_best.pt (8.4 MB)",
    "  SHA256: 19a408c068f19ef917a4fd099656c3c5...",
    "  Incluye: sae_latent_512.pt",
    "  Neuron analysis + AFE/AFC outputs",
], sz=11, cc=LIGHT_GRAY)
section(s, Inches(0.8), Inches(4.8), Inches(11), "Reproducibilidad garantizada")
bullets(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), [
    "Docker multi-stage (python:3.11-slim)",
    "requirements.txt congelado con versiones exactas",
    "Seed fija 42 en todos los experimentos",
    "28 ejecuciones trazadas en runs/ con logs y eventos",
    "Manifest MD5 por archivo para verificacion de integridad",
], sz=12, cc=LIGHT_GRAY)

# === SLIDE 44: ENTREGABLES ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header(s, 42, "Entregables del Proyecto", "Checklist contra la rubrica")
h = ["Componente", "Peso", "Estado", "Detalle"]
r = [
    ["Panel Sit. 1 (50+ GB)", "20%", "COMPLETADO", "66.78 GB, manifest MD5, ETL Dask"],
    ["EDA Panel", "incluido", "COMPLETADO", "5 notebooks, 30+ visualizaciones"],
    ["Modelo CLIP+SAE", "20%", "EN CURSO", "Dataset listo, SAE OK, CLIP mejorando"],
    ["AFE + AFC", "incluido", "COMPLETADO", "KMO=0.97, RMSEA=0.04, CFI=0.999"],
    ["ConvLSTM + Kriging", "30%", "PENDIENTE", "Tensor generado, modelo por entrenar"],
    ["Frontend Next.js", "10%", "AVANZADO", "Mapas, modo oscuro, hooks listos"],
    ["Informe PDF", "10%", "PENDIENTE", "Template LaTeX disponible"],
    ["Pitch y defensa", "10%", "PENDIENTE", "Esta presentacion como base"],
]
tbl(s, Inches(0.4), Inches(1.5), [Inches(2.5), Inches(1.2), Inches(1.5), Inches(5.5)], h, r)
section(s, Inches(0.8), Inches(4.5), Inches(11), "Bonificaciones potenciales")
bullets(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(2.0), [
    "Modo oscuro en frontend: +2 pts  YA IMPLEMENTADO",
    "Tercer modalidad input (audio Whisper): +3 pts",
    "Analisis de equidad espacial por estrato socioeconomico: +4 pts",
    "Comparacion con OMI/AURA como 2da fuente: +3 pts",
], sz=12, cc=LIGHT_GRAY)

# === SLIDE 45: ROADMAP Y CIERRE ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, RGBColor(0x08,0x0E,0x1A))
rect(s, 0, 0, Inches(0.15), H, ACCENT_BLUE)
rect(s, Inches(1.5), Inches(2.3), Inches(10.3), Pt(2), ACCENT_GREEN)
rect(s, Inches(1.5), Inches(5.3), Inches(10.3), Pt(2), ACCENT_BLUE)
tb(s, Inches(1.5), Inches(0.8), Inches(10.5), Inches(0.7), "Roadmap y Cierre", sz=34, cc=WHITE, b=True)

milestones = [
    ("COMPLETADO", "Panel 66.78 GB  ETL Dask  EDA  Manifest MD5", ACCENT_GREEN),
    ("COMPLETADO", "Dataset 2,276 pares  SAE post-train (sparsity 94%)", ACCENT_GREEN),
    ("COMPLETADO", "AFE + AFC: KMO=0.97, RMSEA=0.04, CFI=0.999", ACCENT_GREEN),
    ("COMPLETADO", "Frontend Next.js con mapas, modo oscuro, hooks", ACCENT_GREEN),
    ("EN CURSO", "Entrenamiento CLIP contrastivo (LoRA, 12 bandas)", ACCENT_BLUE),
    ("SIGUIENTE", "Split por escena + ReLU SAE + prompts ES/EN", ACCENT_ORANGE),
    ("SIGUIENTE", "ConvLSTM bidireccional + ST-Kriging", ACCENT_ORANGE),
    ("SIGUIENTE", "LOO-CV DAGMA + Moran I + LISA", ACCENT_RED),
    ("FINAL", "9 mapas de gradiente en app + informe + pitch", ACCENT_GREEN),
]
for i, (label, desc, color) in enumerate(milestones):
    y = Inches(1.5) + i * Inches(0.55)
    circle = rect(s, Inches(0.6), y+Inches(0.05), Inches(0.2), Inches(0.2), color)
    if i < len(milestones)-1:
        rect(s, Inches(0.68), y+Inches(0.25), Pt(3), Inches(0.35), SUBTLE)
    tb(s, Inches(1.1), y, Inches(2.0), Inches(0.3), label, sz=11, cc=color, b=True)
    tb(s, Inches(3.2), y, Inches(9.5), Inches(0.3), desc, sz=10, cc=LIGHT_GRAY)

tb(s, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.5),
   "\"Datos masivos listos  Modelo multimodal en iteracion  Producto visible  Validacion rigurosa en camino\"",
   sz=14, cc=ACCENT_BLUE, al=PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.3),
   "github.com/nicothinn/GeoVision-CLIP-Cali", sz=12, cc=MEDIUM_GRAY, al=PP_ALIGN.CENTER)
tb(s, Inches(4), Inches(6.7), Inches(5.3), Inches(0.4),
   "Gracias  Preguntas y discusion", sz=18, cc=ACCENT_GREEN, b=True, al=PP_ALIGN.CENTER)

# === GUARDAR ===
print(f"Total slides: {len(prs.slides)}")
prs.save(OUTPUT_FILE)
print(f"Presentacion guardada: {OUTPUT_FILE}")
